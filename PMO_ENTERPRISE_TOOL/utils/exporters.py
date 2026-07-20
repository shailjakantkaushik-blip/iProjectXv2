"""Per-tab export helpers — exports look EXACTLY like the on-screen tab.

Pipeline for every figure before export:
  1. _force_labels()    — convert hover-only values into permanent data labels
                          so PNG / PDF / PPT show what users see on mouse-over.
  2. fig_to_png()       — render to PNG via Kaleido at high DPI.

The page is reconstructed as a composite (header + KPI strip + figure grid
+ tables), used as:
  • the PNG download itself,
  • the PDF cover page (with detailed per-figure pages after),
  • the PPTX cover slide (with detailed per-figure slides after).
"""
from __future__ import annotations
import copy
from io import BytesIO
from datetime import datetime
import pandas as pd
import streamlit as st
import plotly.graph_objects as go
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import landscape, A4
from reportlab.pdfgen import canvas as pdfcanvas
from reportlab.lib.utils import ImageReader

from utils.theme_manager import get_theme


# ─────────────────────── helpers ───────────────────────
def _hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    if len(h) == 3: h = "".join(c*2 for c in h)
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def _font(size: int, bold: bool = False):
    paths = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/Library/Fonts/Arial.ttf",
    ]
    for p in paths:
        try: return ImageFont.truetype(p, size)
        except Exception: pass
    return ImageFont.load_default()


def _fmt(v):
    if isinstance(v, (int,)): return f"{v:,}"
    if isinstance(v, float):
        return f"{v:,.0f}" if abs(v) >= 100 else f"{v:,.2f}"
    return str(v)


def _is_empty(v) -> bool:
    """None / '' / empty list / empty numpy array — without triggering
    numpy's 'truth value is ambiguous' on multi-element arrays."""
    if v is None:
        return True
    try:
        return len(v) == 0
    except TypeError:
        return not bool(v)


def _force_labels(fig: go.Figure) -> go.Figure:
    """Make hover values permanently visible by adding on-trace data labels.
    Returns a deep copy — caller's figure is untouched."""
    f = copy.deepcopy(fig)
    for tr in f.data:
        ttype = getattr(tr, "type", "")
        try:
            if ttype == "bar":
                vals = tr.y if (tr.orientation or "v") == "v" else tr.x
                if not _is_empty(vals) and _is_empty(getattr(tr, "text", None)):
                    tr.text = [_fmt(v) for v in vals]
                tr.textposition = tr.textposition or "outside"
                tr.cliponaxis = False
            elif ttype == "scatter":
                # Skip very dense scatters
                yvals = getattr(tr, "y", None)
                n = 0 if _is_empty(yvals) else len(yvals)
                if n and n <= 30:
                    if _is_empty(getattr(tr, "text", None)):
                        tr.text = [_fmt(v) for v in yvals]
                    mode = (tr.mode or "lines+markers")
                    if "text" not in mode: mode = mode + "+text"
                    tr.mode = mode
                    tr.textposition = "top center"
                    tr.textfont = dict(size=10)
            elif ttype == "pie":
                tr.textinfo = "label+percent+value"
                tr.insidetextorientation = "radial"
            elif ttype == "funnel":
                tr.textinfo = "value+percent total"
            elif ttype == "histogram":
                tr.texttemplate = "%{y}"
                tr.textposition = "outside"
            elif ttype == "heatmap":
                # Show numeric values inside heatmap cells
                tr.texttemplate = "%{z:.0f}"
                tr.textfont = dict(size=9)
        except Exception:
            pass
    return f



# ─────────────────────── core: fig → PNG ───────────────────────
def _sanitize_for_json(fig: go.Figure) -> go.Figure:
    """Convert pandas Timestamp / datetime / numpy values inside trace arrays
    AND layout (shapes, annotations, axis range) into JSON-safe equivalents
    so Kaleido's to_image doesn't raise 'Type is not JSON serializable: Timestamp'."""
    import numpy as np
    import datetime as _dt

    def _scalar(v):
        if v is None: return v
        if isinstance(v, pd.Timestamp): return v.isoformat()
        if isinstance(v, np.datetime64): return pd.Timestamp(v).isoformat()
        if isinstance(v, (_dt.datetime, _dt.date)): return v.isoformat()
        if isinstance(v, np.generic):
            try: return v.item()
            except Exception: return str(v)
        return v

    def _conv(seq):
        if seq is None: return seq
        try: return [_scalar(v) for v in seq]
        except TypeError: return seq

    for tr in fig.data:
        for attr in ("x", "y", "z", "text", "base", "customdata"):
            if hasattr(tr, attr):
                val = getattr(tr, attr)
                if val is not None and not isinstance(val, str):
                    try: tr[attr] = _conv(list(val))
                    except Exception: pass

    try:
        for sh in (fig.layout.shapes or []):
            for k in ("x0", "x1", "y0", "y1"):
                v = getattr(sh, k, None)
                if v is not None:
                    try: setattr(sh, k, _scalar(v))
                    except Exception: pass
    except Exception: pass

    try:
        for an in (fig.layout.annotations or []):
            for k in ("x", "y", "ax", "ay"):
                v = getattr(an, k, None)
                if v is not None:
                    try: setattr(an, k, _scalar(v))
                    except Exception: pass
    except Exception: pass

    try:
        for axname in dir(fig.layout):
            if not (axname.startswith("xaxis") or axname.startswith("yaxis")):
                continue
            ax = getattr(fig.layout, axname, None)
            rng = getattr(ax, "range", None) if ax is not None else None
            if not _is_empty(rng) and len(rng) >= 2:
                try: ax.range = [_scalar(rng[0]), _scalar(rng[1])]
                except Exception: pass

    except Exception: pass

    return fig


def fig_to_png(fig: go.Figure, width: int = 1400, height: int = 700) -> bytes:
    safe = _sanitize_for_json(_force_labels(fig))
    return safe.to_image(format="png", width=width, height=height, scale=2)



# ─────────────────────── composite renderer (mirrors the page) ───────────────────────
def _composite(bundle, cols: int = 3, fig_w: int = 880, fig_h: int = 460,
               kpi_h: int = 100, pad: int = 22) -> Image.Image:
    """One big image laid out exactly like the on-screen tab:
        title → KPI strip → figure grid (`cols` per row) → tables."""
    t = get_theme()
    bg     = _hex_to_rgb(t["bg"])
    surf   = _hex_to_rgb(t["surface"])
    text   = _hex_to_rgb(t["text"])
    muted  = _hex_to_rgb(t["muted"])
    border = _hex_to_rgb(t["border"])
    accent = _hex_to_rgb(t["accent"])

    all_figs = bundle.get("figs", [])
    # Pull Governance Flow out so it can be rendered as a full-width hero band
    # ABOVE the regular figure grid (and above the export buttons on screen).
    gov_fig = None
    figs = []
    for n_, f_ in all_figs:
        if gov_fig is None and "governance flow" in str(n_).lower():
            gov_fig = (n_, f_)
        else:
            figs.append((n_, f_))
    kpis   = bundle.get("kpis", [])
    tables = bundle.get("tables", [])

    canvas_w = cols * fig_w + (cols + 1) * pad
    header_h = 110
    kpi_band = kpi_h + pad if kpis else 0
    # Hero governance band (full width)
    gov_h = 0
    if gov_fig is not None:
        gov_h = int(fig_h * 0.95) + pad
    rows     = (len(figs) + cols - 1) // cols if figs else 0
    fig_band = rows * (fig_h + pad) + pad if rows else 0

    # Pre-render tables to estimate height
    f_t = _font(13); f_th = _font(13, True)
    table_blocks = []
    for name, df in tables:
        df_s = df.head(20).astype(str)
        cols_ = list(df_s.columns)[:10]
        rows_n = len(df_s) + 1
        block_h = 30 + rows_n * 20 + 18
        table_blocks.append((name, df_s, cols_, block_h))
    tbl_band = sum(b[3] for b in table_blocks) + (pad if table_blocks else 0)

    total_h = header_h + kpi_band + gov_h + fig_band + tbl_band + pad
    img = Image.new("RGB", (canvas_w, total_h), bg)
    draw = ImageDraw.Draw(img)

    f_title = _font(32, True); f_sub = _font(16); f_lbl = _font(11)
    f_val   = _font(28, True); f_sec = _font(14, True)

    # Header
    draw.text((pad, 22), bundle.get("title", ""), fill=text, font=f_title)
    draw.text((pad, 64), bundle.get("subtitle", ""), fill=muted, font=f_sub)
    draw.line([(pad, header_h-6), (canvas_w-pad, header_h-6)], fill=border, width=1)

    y = header_h
    # KPI strip
    if kpis:
        n = len(kpis)
        card_w = (canvas_w - (n+1)*pad) // n
        for i, (lbl, val) in enumerate(kpis):
            x = pad + i*(card_w+pad)
            draw.rounded_rectangle([x, y, x+card_w, y+kpi_h-pad],
                                   radius=8, outline=border, width=1, fill=surf)
            draw.text((x+12, y+10), str(lbl).upper(), fill=muted, font=f_lbl)
            draw.text((x+12, y+32), str(val), fill=text, font=f_val)
        y += kpi_band

    # Governance flow hero band (full width, above the regular grid)
    if gov_fig is not None:
        name, fig = gov_fig
        band_w = canvas_w - 2 * pad
        band_h = gov_h - pad
        draw.rounded_rectangle([pad-4, y-4, pad+band_w+4, y+band_h+4],
                               radius=10, outline=border, width=1)
        png = fig_to_png(fig, width=band_w*2, height=band_h*2)
        im2 = Image.open(BytesIO(png)).resize((band_w, band_h))
        img.paste(im2, (pad, y))
        y += gov_h

    # Figure grid (3 per row by default)
    for i, (name, fig) in enumerate(figs):
        r, col_ = divmod(i, cols)
        x  = pad + col_ * (fig_w + pad)
        yy = y + r * (fig_h + pad)
        # Frame
        draw.rounded_rectangle([x-4, yy-4, x+fig_w+4, yy+fig_h+4],
                               radius=8, outline=border, width=1)
        png = fig_to_png(fig, width=fig_w*2, height=fig_h*2)
        im2 = Image.open(BytesIO(png)).resize((fig_w, fig_h))
        img.paste(im2, (x, yy))
    y += fig_band

    # Tables
    for name, df_s, cols_, block_h in table_blocks:
        draw.text((pad, y+4), name, fill=text, font=f_sec)
        y += 26
        if not cols_: continue
        cw = (canvas_w - 2*pad) // max(len(cols_), 1)
        # Header background
        draw.rectangle([pad, y, canvas_w-pad, y+22], fill=surf, outline=border)
        for j, c in enumerate(cols_):
            draw.text((pad+j*cw+6, y+4), str(c)[:20], fill=accent, font=f_th)
        y += 22
        for _, row in df_s[cols_].iterrows():
            for j, c in enumerate(cols_):
                draw.text((pad+j*cw+6, y+2), str(row[c])[:24], fill=text, font=f_t)
            draw.line([(pad, y+19), (canvas_w-pad, y+19)], fill=border, width=1)
            y += 20
        y += pad
    return img


# ─────────────────────── PNG ───────────────────────
def bundle_to_png(bundle) -> bytes:
    out = BytesIO()
    _composite(bundle).save(out, "PNG", optimize=True)
    return out.getvalue()


# ─────────────────────── PDF ───────────────────────
def bundle_to_pdf(bundle) -> bytes:
    """Page 1 = composite replica. Pages 2..N = one figure per page (zoomed
    for legibility) + tables."""
    buf = BytesIO()
    pagesize = landscape(A4); W, H = pagesize
    c = pdfcanvas.Canvas(buf, pagesize=pagesize)
    t = get_theme()
    bg = tuple(v/255 for v in _hex_to_rgb(t["bg"]))
    fg = tuple(v/255 for v in _hex_to_rgb(t["text"]))
    mut = tuple(v/255 for v in _hex_to_rgb(t["muted"]))
    accent = tuple(v/255 for v in _hex_to_rgb(t["accent"]))

    def _fill_bg():
        c.setFillColorRGB(*bg); c.rect(0, 0, W, H, fill=1, stroke=0)

    # ── Page 1: composite snapshot of the tab ──
    _fill_bg()
    composite = _composite(bundle)
    cbuf = BytesIO(); composite.save(cbuf, "PNG"); cbuf.seek(0)
    img = ImageReader(cbuf); iw, ih = img.getSize()
    s = min((W-40)/iw, (H-40)/ih)
    c.drawImage(img, (W-iw*s)/2, (H-ih*s)/2, width=iw*s, height=ih*s)
    c.showPage()

    # ── per-figure detail pages ──
    for name, fig in bundle.get("figs", []):
        _fill_bg()
        c.setFillColorRGB(*fg); c.setFont("Helvetica-Bold", 16)
        c.drawString(36, H-40, f"{bundle.get('title','')} · {name}")
        png = fig_to_png(fig, width=1800, height=1000)
        img = ImageReader(BytesIO(png)); iw, ih = img.getSize()
        s = min((W-72)/iw, (H-90)/ih)
        c.drawImage(img, (W-iw*s)/2, 30, width=iw*s, height=ih*s)
        c.showPage()

    # ── table pages ──
    for name, df in bundle.get("tables", []):
        _fill_bg()
        c.setFillColorRGB(*fg); c.setFont("Helvetica-Bold", 16)
        c.drawString(36, H-40, name)
        c.setFont("Helvetica", 8)
        df_s = df.head(40).astype(str)
        cols = list(df_s.columns)[:10]
        col_w = (W-72)/max(len(cols), 1)
        y = H-70
        c.setFillColorRGB(*accent)
        for i, col in enumerate(cols):
            c.drawString(36 + i*col_w, y, str(col)[:18])
        c.setFillColorRGB(*fg)
        y -= 14
        for _, row in df_s[cols].iterrows():
            for i, col in enumerate(cols):
                c.drawString(36 + i*col_w, y, str(row[col])[:22])
            y -= 11
            if y < 30: break
        c.showPage()

    c.save(); return buf.getvalue()


# ─────────────────────── PPTX ───────────────────────
def _slide_bg(slide, rgb):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = rgb


def _title(slide, txt, sub, fg, sub_color):
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(1))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = fg
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub
        p2.font.size = Pt(12); p2.font.color.rgb = sub_color


def _add_table(slide, df, left, top, width, height, fg, max_rows=14, max_cols=8):
    df = df.head(max_rows).iloc[:, :max_cols].copy()
    rows, cols = df.shape[0]+1, df.shape[1]
    if cols == 0: return
    tbl = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                 Inches(width), Inches(height)).table
    for j, c_ in enumerate(df.columns):
        cell = tbl.cell(0, j); cell.text = str(c_)
        for p in cell.text_frame.paragraphs:
            for r in p.runs: r.font.size=Pt(9); r.font.bold=True; r.font.color.rgb=fg
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j); cell.text = str(v)[:40]
            for p in cell.text_frame.paragraphs:
                for r in p.runs: r.font.size=Pt(8); r.font.color.rgb=fg


def bundle_to_pptx(bundles, single_tab=False) -> bytes:
    if single_tab: bundles = [bundles]
    t = get_theme()
    bg_rgb  = RGBColor(*_hex_to_rgb(t["bg"]))
    fg_rgb  = RGBColor(*_hex_to_rgb(t["text"]))
    sub_rgb = RGBColor(*_hex_to_rgb(t["muted"]))

    prs = Presentation()
    prs.slide_width  = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    if not single_tab:
        s = prs.slides.add_slide(blank); _slide_bg(s, bg_rgb)
        _title(s, "PMO Portfolio — Executive Briefing",
               datetime.now().strftime("%d %b %Y"), fg_rgb, sub_rgb)

    for bundle in bundles:
        # ── COVER slide = composite snapshot (mirrors the on-screen tab) ──
        s = prs.slides.add_slide(blank); _slide_bg(s, bg_rgb)
        composite = _composite(bundle)
        cbuf = BytesIO(); composite.save(cbuf, "PNG"); cbuf.seek(0)
        iw, ih = composite.size
        # Fit within slide (13.33 × 7.5 inches)
        slide_w_in, slide_h_in = 13.0, 7.2
        ratio = min(slide_w_in / (iw / 96), slide_h_in / (ih / 96))
        w_in = (iw / 96) * ratio; h_in = (ih / 96) * ratio
        left = (13.333 - w_in) / 2; top = (7.5 - h_in) / 2
        s.shapes.add_picture(cbuf, Inches(left), Inches(top),
                             width=Inches(w_in), height=Inches(h_in))

        # ── per-figure detail slides (with labels forced on) ──
        for name, fig in bundle.get("figs", []):
            s = prs.slides.add_slide(blank); _slide_bg(s, bg_rgb)
            _title(s, f"{bundle.get('title','')} · {name}", "", fg_rgb, sub_rgb)
            png = fig_to_png(fig, width=1800, height=1000)
            s.shapes.add_picture(BytesIO(png), Inches(0.4), Inches(1.2),
                                 width=Inches(12.5))

        # ── table slides ──
        for name, df in bundle.get("tables", [])[:2]:
            s = prs.slides.add_slide(blank); _slide_bg(s, bg_rgb)
            _title(s, f"{bundle.get('title','')} · {name}", "", fg_rgb, sub_rgb)
            _add_table(s, df, 0.4, 1.4, 12.5, 5.5, fg_rgb)

    out = BytesIO(); prs.save(out); return out.getvalue()


# ─────────────────────── Exec Dashboard PDF (with all timelines expanded) ───────────────────────
def executive_dashboard_pdf(bundle, timeline_figs, view_label: str = "") -> bytes:
    """PDF that exactly mirrors the Executive Dashboard on screen, followed by
    one page per project-timeline group with ALL projects expanded.

    Parameters
    ----------
    bundle        : the same bundle passed to `render_export_buttons` (KPIs,
                    figures, tables, title).
    timeline_figs : list of `(group_label, plotly_figure)` from
                    `timeline_views.build_bucket_timeline_figs(...)`.
    view_label    : e.g. "Portfolio View" — printed on the section divider.
    """
    buf = BytesIO()
    pagesize = landscape(A4); W, H = pagesize
    c = pdfcanvas.Canvas(buf, pagesize=pagesize)
    t = get_theme()
    bg  = tuple(v/255 for v in _hex_to_rgb(t["bg"]))
    fg  = tuple(v/255 for v in _hex_to_rgb(t["text"]))
    mut = tuple(v/255 for v in _hex_to_rgb(t["muted"]))

    def _fill_bg():
        c.setFillColorRGB(*bg); c.rect(0, 0, W, H, fill=1, stroke=0)

    # Pages 1..N — exact dashboard composite, sliced across multiple pages
    # so every KPI / chart / table renders at readable size (no squish-to-fit).
    composite = _composite(bundle)
    iw, ih = composite.size
    margin = 24
    avail_w = W - 2 * margin
    avail_h = H - 2 * margin
    scale = avail_w / iw                       # fit width, keep aspect
    scaled_h = ih * scale
    # Number of pages needed to show the full composite at that scale
    n_pages = max(1, int((scaled_h + avail_h - 1) // avail_h))
    slice_src_h = int(avail_h / scale)         # source-pixel height per page
    for p in range(n_pages):
        _fill_bg()
        y0 = p * slice_src_h
        y1 = min(ih, y0 + slice_src_h)
        strip = composite.crop((0, y0, iw, y1))
        sbuf = BytesIO(); strip.save(sbuf, "PNG"); sbuf.seek(0)
        simg = ImageReader(sbuf); siw, sih = simg.getSize()
        draw_w = siw * scale
        draw_h = sih * scale
        c.drawImage(simg, margin, H - margin - draw_h,
                    width=draw_w, height=draw_h)
        # tiny page indicator
        c.setFillColorRGB(*mut); c.setFont("Helvetica", 8)
        c.drawRightString(W - margin, 12,
                          f"Executive Dashboard — page {p+1} of {n_pages}")
        c.showPage()


    # Divider page
    if timeline_figs:
        _fill_bg()
        c.setFillColorRGB(*fg); c.setFont("Helvetica-Bold", 28)
        c.drawString(40, H-80, "Portfolio Timelines — All Projects Expanded")
        c.setFillColorRGB(*mut); c.setFont("Helvetica", 14)
        c.drawString(40, H-110, f"Grouped by: {view_label or '—'}   ·   "
                                f"{len(timeline_figs)} group(s)   ·   "
                                f"{datetime.now().strftime('%d %b %Y %H:%M')}")
        c.showPage()

    # One page per timeline group; scale to fit
    for label, fig in timeline_figs:
        _fill_bg()
        c.setFillColorRGB(*fg); c.setFont("Helvetica-Bold", 16)
        c.drawString(36, H-36, f"{label}")
        try:
            # Use the figure's own height as a hint for aspect ratio
            fig_h_hint = int(getattr(fig.layout, "height", 700) or 700)
            png = fig_to_png(fig, width=1800, height=max(700, min(fig_h_hint, 2600)))
        except Exception as e:
            c.setFont("Helvetica", 10); c.setFillColorRGB(1, 0, 0)
            c.drawString(36, H-60, f"Render failed: {e}")
            c.showPage(); continue
        img = ImageReader(BytesIO(png)); iw, ih = img.getSize()
        s = min((W-72)/iw, (H-80)/ih)
        c.drawImage(img, (W-iw*s)/2, 30, width=iw*s, height=ih*s)
        c.showPage()

    c.save(); return buf.getvalue()




# ─────────────────────── Streamlit helper ───────────────────────
def render_export_buttons(bundle):
    st.markdown("---")
    st.markdown("#### 📤 Export this tab (PNG / PDF / PPT — exact replica)")
    st.caption("Exports include hover values as permanent data labels and "
               "mirror the on-screen layout.")
    c1, c2, c3 = st.columns(3)
    slug = bundle.get("title", "tab").split("·")[0].strip().replace(" ", "") + \
           "_" + datetime.now().strftime("%Y%m%d_%H%M%S")
    try:
        png = bundle_to_png(bundle)
        c1.download_button("🖼  PNG Image", png, file_name=f"{slug}.png",
                           mime="image/png", use_container_width=True)
    except Exception as e:
        c1.error(f"PNG: {e}")
    try:
        pdf = bundle_to_pdf(bundle)
        c2.download_button("📄 PDF", pdf, file_name=f"{slug}.pdf",
                           mime="application/pdf", use_container_width=True)
    except Exception as e:
        c2.error(f"PDF: {e}")
    try:
        pptx = bundle_to_pptx(bundle, single_tab=True)
        c3.download_button("📊 PowerPoint", pptx, file_name=f"{slug}.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                           use_container_width=True)
    except Exception as e:
        c3.error(f"PPTX: {e}")
